from fastapi import Body
from pptx.util import Inches, Pt
import logging
import io
import os
from testdevices import TranslatorModel
from typing import List
from fastapi import FastAPI
from fastapi import FastAPI, Body
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from docx import Document
from pydantic import BaseModel
import tkinter as tk
from tkinter import filedialog
import os
import subprocess
import platform
import time


app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- グローバル変数 ---
prs = None
filepath = None

# --- Pydantic Models ---


class TextToTranslate(BaseModel):
    text: str


class ShapeUpdate(BaseModel):
    shape_index: int
    translated_text: str


class SlideUpdateShapes(BaseModel):
    slide_index: int
    shapes: list[ShapeUpdate]


class SlidesPayload(BaseModel):
    slides: list[SlideUpdateShapes]

# --- Util ---


def get_color(font_color):
    """フォントの色 (RGB) を文字列で取得します。"""
    try:
        if font_color and font_color.rgb:
            return str(font_color.rgb)
    except:
        return None
    return None


def close_all_powerpoint_presentations_mac(save_changes=False):
    """
    🚨 Mac専用機能 🚨
    Mac上でAppleScriptを使って、現在開いている全てのPowerPointプレゼンテーションを閉じます。

    Args:
        save_changes (bool): 
            True: 変更を保存してから閉じます。
            False: 保存せずに閉じます（未保存の変更は破棄されます）。
    """
    if platform.system() != "Darwin":
        print("⚠️ 警告: Mac環境ではないため、PowerPointを閉じる操作はスキップされました。")
        return

    if save_changes:
        # 変更を保存して閉じるためのAppleScript
        script_command = """
        tell application "Microsoft PowerPoint"
            # 開いているプレゼンテーションを全て閉じる (保存して閉じる)
            close every presentation saving yes
        end tell
        """
        print("✨ [Mac] PowerPointの全てのプレゼンテーションを、保存して閉じます...")
    else:
        # 変更を保存せずに閉じるためのAppleScript (savedプロパティをtrueに設定して強制的に閉じる)
        script_command = """
        tell application "Microsoft PowerPoint"
            try
                set allPresentations to presentations
                
                repeat with i from (count of allPresentations) to 1 by -1
                    set aPresentation to item i of allPresentations
                    
                    # 変更がされていても保存を促さないようにSavedプロパティをtrueに設定
                    set saved of aPresentation to true
                    
                    # プレゼンテーションを閉じる
                    close aPresentation
                end repeat
            on error errMsg
                log "AppleScriptエラー: " & errMsg
            end try
        end tell
        """
        print("✨ [Mac] PowerPointの全てのプレゼンテーションを、保存せずに閉じます...")

    try:
        # PythonからAppleScriptを実行します
        subprocess.run(['osascript', '-e', script_command],
                       check=True, capture_output=True, text=True)
        print("✅ [Mac] PowerPointの閉じる処理が完了しました。")

    except subprocess.CalledProcessError as e:
        print(f"🚨 [Mac] AppleScriptの実行中にエラーが発生しました: {e.stderr.strip()}")
    except FileNotFoundError:
        # osascriptが見つからないのは通常ありえませんが、念のため
        print("🚨 'osascript' コマンドが見つかりません。")


@app.get("/wait")
def wait():
    time.sleep(3)  # ← 3秒待つ
    return {"status": "ok", "message": "API response arrived!"}

# ----------------------------------------------------
# /get_file
# ----------------------------------------------------


@app.get("/get_file")
async def load_file():
    global prs, filepath

    # tkinterはGUIアプリなので、サーバー環境によっては非推奨
    root = tk.Tk()
    root.withdraw()

    path = filedialog.askopenfilename(
        title="ファイルを選択",
        filetypes=[("PowerPoint", "*.pptx"),
                   ("Word", "*.docx")]
    )

    if not path:
        return {"error": "ファイルが選択されていません"}

    filepath = path
    print("選択されたファイルパス:", filepath)
    filename = os.path.basename(path)
    ext = os.path.splitext(path)[1].lower()
    slides = []

    # ---------------------------------------
    # PPTX 読み込み
    # ---------------------------------------
    if ext == ".pptx":
        prs = Presentation(path)

        # デスクトップでファイルを開く
        if platform.system() == "Darwin":
            subprocess.run(["open", path])
        elif platform.system() == "Windows":
            os.startfile(path)

        # データ抽出
        for i, slide in enumerate(prs.slides):
            slide_shapes = []

            for shape_index, shape in enumerate(slide.shapes):
                shape_data = {
                    "shape_index": shape_index,
                    "left": shape.left if hasattr(shape, "left") else None,
                    "top": shape.top if hasattr(shape, "top") else None,
                    "width": shape.width if hasattr(shape, "width") else None,
                    "height": shape.height if hasattr(shape, "height") else None,
                    "paragraphs": []
                }

                if shape.has_text_frame:
                    tf = shape.text_frame

                    for p_index, paragraph in enumerate(tf.paragraphs):
                        paragraph_data = {
                            "paragraph_index": p_index,
                            "text": paragraph.text,
                            "runs": []
                        }

                        for r_index, run in enumerate(paragraph.runs):
                            paragraph_data["runs"].append({
                                "run_index": r_index,
                                "text": run.text,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "size": run.font.size.pt if run.font.size else None,
                                "color": get_color(run.font.color)
                            })

                        shape_data["paragraphs"].append(paragraph_data)

                slide_shapes.append(shape_data)

            slides.append({"index": i, "shapes": slide_shapes})

    # ---------------------------------------
    # DOCX 読み込み（200文字チャンク＋段落分割）
    # ---------------------------------------
    elif ext == ".docx":
        prs = None
        doc = Document(path)

        # 全段落取得（空段落除外）
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        slides = []
        slide_index = 0
        chunk_size = 200
        current_chunk_text = ""
        current_paragraphs = []

        for paragraph in paragraphs:
            while paragraph:
                remaining_space = chunk_size - len(current_chunk_text)
                if len(paragraph) <= remaining_space:
                    # チャンクに追加
                    current_paragraphs.append({
                        "paragraph_index": len(current_paragraphs),
                        "text": paragraph,
                        "runs": []
                    })
                    current_chunk_text += paragraph
                    paragraph = ""
                else:
                    # チャンクに入りきらない場合、分割
                    part = paragraph[:remaining_space]
                    current_paragraphs.append({
                        "paragraph_index": len(current_paragraphs),
                        "text": part,
                        "runs": []
                    })
                    current_chunk_text += part
                    paragraph = paragraph[remaining_space:]

                # チャンクがいっぱいになったらスライド化
                if len(current_chunk_text) >= chunk_size:
                    slides.append({
                        "index": slide_index,
                        "shapes": [
                            {
                                "shape_index": 0,
                                "text": current_chunk_text,
                                "paragraphs": current_paragraphs
                            }
                        ]
                    })
                    slide_index += 1
                    current_chunk_text = ""
                    current_paragraphs = []

        # 残りがあれば最後のスライドに追加
        if current_chunk_text:
            slides.append({
                "index": slide_index,
                "shapes": [
                    {
                        "shape_index": 0,
                        "text": current_chunk_text,
                        "paragraphs": current_paragraphs
                    }
                ]
            })

    else:
        return {"error": f"{ext}形式は未対応です"}

    return {
        "path": path,
        "filepath": filepath,
        "filename": filename,
        "slides": slides,
        "ext": ext
    }

# ----------------------------------------------------
# /translate_text (翻訳モデルがないため、このAPIは機能しない可能性があります)
# ----------------------------------------------------


# モデルをグローバルロード
TRANS_MODEL = TranslatorModel(model_dir="openvino_model")

# データモデル定義


class Paragraph(BaseModel):
    text: str


class Shape(BaseModel):
    paragraphs: List[Paragraph]


class Slide(BaseModel):
    shapes: List[Shape]


class TextsToTranslate(BaseModel):
    texts: List[str]
    language: str


@app.post("/translate_texts")
async def api_translate_texts(data:    TextsToTranslate):
    if TRANS_MODEL is None:
        return {"error": "翻訳モデルがロードされていません", "translated_texts": []}

    tgt_lang = "ja_XX" if data.language == "ja" else "en_XX"
    src_lang = "en_XX" if tgt_lang == "ja_XX" else "ja_XX"

    translated_texts = []

    for text in data.texts:
        print("Translating text:", text)
        # 翻訳する
        translated_text = TRANS_MODEL.translate_text(
            text,
            src_lang=src_lang,
            tgt_lang=tgt_lang
        )
        translated_texts.append(translated_text)

    # 翻訳されたテキストのリストを返す
    return {"status": "ok", "translated_texts": translated_texts}


# from fastapi import FastAPI
# from pydantic import BaseModel
# from typing import List, Optional
# import openvino as ov
# from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
# import numpy as np

# # ---------------------------
# # FastAPI 用データモデル
# # ---------------------------
# class Paragraph(BaseModel):
#     text: str

# class Shape(BaseModel):
#     paragraphs: Optional[List[Paragraph]] = []

# class Slide(BaseModel):
#     shapes: Optional[List[Shape]] = []

# class FileData(BaseModel):
#     slides: Optional[List[Slide]] = []

# # ---------------------------
# # OpenVINO 翻訳モデル初期化（アプリ起動時に一度だけ）
# # ---------------------------
# MODEL_NAME = "Helsinki-NLP/opus-mt-ja-en"

# tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME)
# pt_model = AutoModelForSeq2SeqLM.from_pretrained(MODEL_NAME).eval()

# example_input_ids = tokenizer("こんにちは", return_tensors="pt").input_ids
# ov_model = ov.convert_model(pt_model, example_input=(example_input_ids,))

# core = ov.Core()
# compiled_model = core.compile_model(ov_model, device_name="GPU")  # GPU使用

# # ---------------------------
# # 翻訳関数（OpenVINO 用）
# # ---------------------------
# def ov_translate_text(text: str, max_tokens=50) -> str:
#     """
#     OpenVINO を使って文章を日本語→英語に翻訳する関数
#     """
#     if not text.strip():
#         return text

#     input_ids = tokenizer(text, return_tensors="np").input_ids  # numpy 配列

#     for _ in range(max_tokens):
#         logits = compiled_model({"input_ids": input_ids})[compiled_model.output(0)]
#         next_token = np.argmax(logits[0, -1, :])
#         input_ids = np.concatenate([input_ids, [[next_token]]], axis=1)
#         if next_token == tokenizer.eos_token_id:
#             break

#     return tokenizer.decode(input_ids[0], skip_special_tokens=True)


# @app.post("/update_slide")
# def update_slide(data: SlidesPayload):
#     global prs, filepath
#     if prs is None:
#         # prsがNoneの場合、ファイルを再ロードする（グローバル変数の状態がリセットされた場合を想定）
#         prs = Presentation(filepath)

#     for slide_item in data.slides:
#         slide = prs.slides[slide_item.slide_index]
#         for shape_item in slide_item.shapes:
#             shape = slide.shapes[shape_item.shape_index]
#             if shape.has_text_frame:
#                 # この実装はシェイプ全体のテキストを置き換えます
#                 shape.text = shape_item.translated_text

#     # update_slideでは保存しない（savefileでまとめて保存する想定）
#     # prs.save(filepath)
#     return {"status": "ok"}

# ----------------------------------------------------
# /test (シェイプの座標取得)
# ----------------------------------------------------


@app.post("/savedocx")
def save_docx_endpoint(payload: dict = Body(...)):
    """
    DOCXのチャンク翻訳結果を上書き保存するエンドポイント。
    payload = {
        "selectedFilePath": "xxx.docx",
        "chunks": [ "翻訳後テキスト…", "翻訳後テキスト…", ... ]
    }
    """
    selectedFilePath = payload.get("selectedFilePath")
    chunks = payload.get("chunks")

    if not selectedFilePath:
        return {"status": "error", "message": "File path missing"}

    if chunks is None:
        return {"status": "error", "message": "No translated chunks provided"}

    # ---- 元ファイル読み込み ----
    try:
        doc = Document(selectedFilePath)
    except Exception as e:
        return {"status": "error", "message": f"Cannot open DOCX: {str(e)}"}

    # ---- 翻訳済みチャンク → 1本のテキストに結合（改行そのまま）----
    merged_text = "\n".join(chunks)

    # ---- DOCX の段落を全て削除 ----
    for _ in range(len(doc.paragraphs)):
        p = doc.paragraphs[0]
        p._element.getparent().remove(p._element)

    # ---- 新しい段落として書き込み（改行保持）----
    for line in merged_text.split("\n"):
        p = doc.add_paragraph()
        p.add_run(line)

    # ---- 上書き保存 ----
    try:
        doc.save(selectedFilePath)
    except Exception as e:
        return {"status": "error", "message": f"Failed to save DOCX: {str(e)}"}

    return {
        "status": "ok",
        "saved_path": selectedFilePath
    }


# ----------------------------------------------------
# /savefile (ユーザーの要望通り、閉じる処理を残す)
# ----------------------------------------------------
logging.basicConfig(level=logging.INFO)


@app.post("/saveppt")
def save_ppt_endpoint(payload: dict = Body(...)):
    """
    フロントエンドから送られた単一スライドのシェイプデータをPPTXファイルに保存（上書き）する。
    """
    selectedFilePath = payload.get("selectedFilePath")
    print(payload)
    slide_index_to_update = payload.get("slide_index")
    shapes_data = payload.get("shapes", [])
    print("★★ selectedFilePath:", selectedFilePath)

    logging.info(f"Selected file path: {selectedFilePath}")
    logging.info(f"Slide index to update: {slide_index_to_update}")

    # --- 1. エラーチェック ---
    if not selectedFilePath:
        return {"status": "error", "message": "File path is missing"}

    if slide_index_to_update is None:
        return {"status": "error", "message": "Slide index is missing"}

    # --- 2. ファイルの読み込み ---
    try:
        with open(selectedFilePath, "rb") as f:
            pptx_bytes = f.read()
    except FileNotFoundError:
        return {"status": "error", "message": f"File not found: {selectedFilePath}"}

    prs = Presentation(io.BytesIO(pptx_bytes))

    # --- 3. 特定のスライドのテキストを更新 ---
    try:
        slide = prs.slides[slide_index_to_update]
    except IndexError:
        return {
            "status": "error",
            "message": f"Invalid slide index: {slide_index_to_update}. Total slides: {len(prs.slides)}"
        }

    for shape_data in shapes_data:
        shape_index = shape_data.get("shape_index")
        text_content = shape_data.get("text", "")

        if shape_index is None:
            logging.warning(
                "Received shape data without shape_index. Skipping.")
            continue

        try:
            shape = slide.shapes[shape_index]

            if hasattr(shape, "text_frame"):
                shape.text = text_content
                logging.info(
                    f"Updated slide {slide_index_to_update}, shape {shape_index}")
            else:
                logging.info(
                    f"Shape {shape_index} on slide {slide_index_to_update} is not a text shape. Skipping.")

        except IndexError:
            logging.warning(
                f"Shape index {shape_index} not found on slide {slide_index_to_update}. Skipping.")
            continue

    # --- 4. 上書き保存 ---
    save_path = selectedFilePath  # 別ファイルを作らず上書き保存

    try:
        prs.save(save_path)
    except Exception as e:
        logging.error(f"Failed to save the presentation: {str(e)}")
        return {"status": "error", "message": f"Failed to save: {str(e)}"}

    logging.info(f"File saved at: {save_path}")

    return {"status": "ok", "saved_path": save_path}


class TextData(BaseModel):
    text: str
    left: float
    top: float
    width: float
    height: float


@app.post("/insert")
def insert_slide(data: TextData):
    global prs, filepath

    try:
        # 💡 修正点2: tryブロック内の処理を正しくインデント
        # print("インサートのパス",filepath)
        # print(data.left,data.top,data.width,data.height)

        slide = prs.slides[0]

        # 座標とサイズを Inches で指定 (例として左上から2インチ、幅4インチなど)
        left = Inches(data.left)
        top = Inches(data.top)
        width = Inches(data.width)
        height = Inches(data.height)

        # left = Inches(7)
        # top = Inches(5)
        # width = Inches(4)
        # height = Inches(5)

        # 指定した座標とサイズでテキストボックスを追加
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        # 受信したテキストを挿入
        p = tf.paragraphs[0]
        p.text = data.text
        # ファイルが見つかった場合の処理をここに続ける
        print(f"ファイル 'input.pptx' を開きました。")
        # print(f"受け取ったテキスト: {data.text}")

        prs.save(filepath)
        print(f"テキストボックスを追加し、ファイルを保存しました。")

        # 実際のPPTX処理 (例: スライドにテキストを追加するコードなど) はここに追加します
        # 例: slide = prs.slides[0]; ...

    except FileNotFoundError:
        # 💡 修正点3: exceptブロック内の処理を正しくインデント
        print("input.pptx が見つかりません。ファイルを作成してください。")

        # ファイルが見つからない場合は、エラーメッセージを返して終了する
        # FastAPIでは exit() ではなく、適切なエラーレスポンスを返すのが一般的です。
        return {"status": "error", "message": "処理に必要な 'input.pptx' ファイルが見つかりませんでした。"}

    # try...exceptブロックの外に出すことで、エラーが発生しなかった場合のみ実行される
    return {"status": "ok", "message": "新しいスライドを追加しました。"}


class TextData(BaseModel):
    text: str


# DOCX ファイルパス（例: input.docx）
docx_filepath = filepath


@app.post("/insert-docx")
def insert_docx(data: TextData):
    try:
        # ファイルが存在しない場合は新規作成
        if os.path.exists(docx_filepath):
            doc = Document(docx_filepath)
            print(f"ファイル '{docx_filepath}' を開きました。")
        else:
            doc = Document()
            print(f"ファイル '{docx_filepath}' が見つからないため新規作成しました。")

        # 新しい段落としてテキストを追加
        doc.add_paragraph(data.text)
        print(f"テキストを追加しました: {data.text}")

        # 保存
        doc.save(docx_filepath)
        print(f"ファイル '{docx_filepath}' に保存しました。")

    except Exception as e:
        print("error:", e)
        return {"status": "error", "message": f"DOCXへの挿入中にエラーが発生しました: {e}"}

    return {"status": "ok", "message": "DOCXにテキストを追加しました。"}


class TranslateOnly(BaseModel):
    text: str
    target_language: str  # "ja" か "en"


@app.post("/insert-translate")
def insert_and_translate(data: TranslateOnly):
    try:
        # 翻訳方向を決定
        tgt_lang = "ja_XX" if data.target_language == "ja" else "en_XX"
        src_lang = "en_XX" if tgt_lang == "ja_XX" else "ja_XX"

        # 翻訳
        translated_text = TRANS_MODEL.translate_text(
            data.text,
            src_lang=src_lang,
            tgt_lang=tgt_lang
        )

        return {
            "status": "ok",
            "translated_text": translated_text
        }

    except Exception as e:
        print("error:", e)
        return {"status": "error", "message": "翻訳失敗"}


class TranslateOnly(BaseModel):
    text: str

# ---------------------
# 🔵 DOCX翻訳専用エンドポイント
# ---------------------


# @app.post("/insert-translate-docx")
# def insert_and_translate_docx(data: TranslateOnly):
#     try:
#         # ① 翻訳
#         translated_text = TRANS_MODEL.translate_text(data.text)

#         # ② DOCXに書き込む場合のサンプル（任意）
#         # doc = Document()
#         # doc.add_paragraph(translated_text)
#         # buffer = io.BytesIO()
#         # doc.save(buffer)
#         # buffer.seek(0)

#         return {
#             "status": "ok",
#             "translated_text": translated_text
#             # "docx_file": buffer.getvalue()  # 必要に応じてバイト配列を返せる
#         }

#     except Exception as e:
#         print("error:", e)
#         return {"status": "error", "message": "翻訳失敗"}


# class ShapeItem(BaseModel):
#     shape_index: int
#     text: str

# class SavePayload(BaseModel):
#     slide_index: int
#     shapes: list[ShapeItem]


# @app.post("/saveppt")
# def saveppt(data: SavePayload):
#     global prs, filepath

#     slide = prs.slides[data.slide_index]

#     for item in data.shapes:
#         try:
#             shape = slide.shapes[item.shape_index]
#         except:
#             continue

#         if not shape.has_text_frame:
#             continue

#         tf = shape.text_frame
#         tf.clear()
#         tf.text = item.text

#     prs.save(filepath)

#     return {"status": "ok"}
